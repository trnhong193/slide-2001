#!/usr/bin/env python3
"""
Resize PowerPoint slide dimensions to match output slide size.

Usage:
    python resize_slide.py <input_pptx> [output_pptx]

This script resizes all slides in a presentation to:
- Width: 720pt (10 inches)
- Height: 405pt (5.625 inches) - 16:9 aspect ratio
"""

import sys
from pathlib import Path

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    HAS_PPTX = True
except ImportError:
    HAS_PPTX = False
    print("Error: python-pptx not installed. Install with: pip install python-pptx")
    sys.exit(1)

# Output slide dimensions (from generate_from_json.js)
SLIDE_WIDTH_PT = 720   # pt
SLIDE_HEIGHT_PT = 405  # pt

# Convert to inches and EMU
# 1 inch = 72 points
SLIDE_WIDTH_IN = SLIDE_WIDTH_PT / 72.0      # 10 inches
SLIDE_HEIGHT_IN = SLIDE_HEIGHT_PT / 72.0    # 5.625 inches

# 1 inch = 914400 EMU (English Metric Units)
SLIDE_WIDTH_EMU = int(SLIDE_WIDTH_IN * 914400)   # 9144000 EMU
SLIDE_HEIGHT_EMU = int(SLIDE_HEIGHT_IN * 914400) # 5143500 EMU


def scale_shape(shape, scale_x, scale_y):
    """
    Scale a shape and all its elements.
    
    Args:
        shape: Shape object from python-pptx
        scale_x: Horizontal scale factor
        scale_y: Vertical scale factor
    """
    try:
        # Scale position and size
        shape.left = int(shape.left * scale_x)
        shape.top = int(shape.top * scale_y)
        shape.width = int(shape.width * scale_x)
        shape.height = int(shape.height * scale_y)
        
        # Scale text font sizes if shape has text
        if hasattr(shape, 'text_frame'):
            try:
                for paragraph in shape.text_frame.paragraphs:
                    if paragraph.font and paragraph.font.size:
                        # Scale font size (use average of scale_x and scale_y for proportional scaling)
                        avg_scale = (scale_x + scale_y) / 2.0
                        paragraph.font.size = Pt(int(paragraph.font.size.pt * avg_scale))
                    
                    # Scale font sizes in runs
                    for run in paragraph.runs:
                        if run.font and run.font.size:
                            avg_scale = (scale_x + scale_y) / 2.0
                            run.font.size = Pt(int(run.font.size.pt * avg_scale))
            except Exception as e:
                # Some shapes might not support font scaling
                pass
        
        # Handle group shapes (recursively scale grouped shapes)
        if shape.shape_type == 6:  # MSO_SHAPE_TYPE.GROUP = 6
            try:
                for grouped_shape in shape.shapes:
                    scale_shape(grouped_shape, scale_x, scale_y)
            except Exception as e:
                pass
        
        # Handle table shapes
        if hasattr(shape, 'table'):
            try:
                # Table cells are already scaled with the table shape
                # But we might need to scale font sizes in cells
                for row in shape.table.rows:
                    for cell in row.cells:
                        if hasattr(cell, 'text_frame'):
                            for paragraph in cell.text_frame.paragraphs:
                                if paragraph.font and paragraph.font.size:
                                    avg_scale = (scale_x + scale_y) / 2.0
                                    paragraph.font.size = Pt(int(paragraph.font.size.pt * avg_scale))
                                for run in paragraph.runs:
                                    if run.font and run.font.size:
                                        avg_scale = (scale_x + scale_y) / 2.0
                                        run.font.size = Pt(int(run.font.size.pt * avg_scale))
            except Exception as e:
                pass
                
    except Exception as e:
        # Skip shapes that can't be scaled
        pass


def resize_presentation(input_path, output_path=None):
    """
    Resize all slides in a PowerPoint presentation to match output slide dimensions.
    Also scales all shapes, images, and text elements proportionally.
    
    Args:
        input_path: Path to input PowerPoint file
        output_path: Path to output file (default: overwrite input)
    """
    input_file = Path(input_path).resolve()
    
    if not input_file.exists():
        print(f"Error: File not found: {input_file}")
        return False
    
    if output_path is None:
        output_path = input_file
    else:
        output_path = Path(output_path).resolve()
    
    print(f"Loading presentation: {input_file}")
    pres = Presentation(str(input_file))
    
    # Get current dimensions
    current_width_emu = pres.slide_width
    current_height_emu = pres.slide_height
    current_width_in = current_width_emu / 914400.0
    current_height_in = current_height_emu / 914400.0
    
    print(f"Current slide dimensions: {current_width_in:.2f}\" x {current_height_in:.2f}\"")
    print(f"Current slide dimensions: {current_width_emu} EMU x {current_height_emu} EMU")
    
    # Check if shapes are already out of bounds (indicating previous resize without scaling)
    # If shapes extend beyond current slide dimensions, we need to calculate original dimensions
    shapes_need_scaling = False
    slides_list = list(pres.slides)
    for i, slide in enumerate(slides_list[:min(3, len(slides_list))]):  # Check first 3 slides
        for shape in slide.shapes:
            if (shape.left + shape.width > pres.slide_width or 
                shape.top + shape.height > pres.slide_height):
                shapes_need_scaling = True
                break
        if shapes_need_scaling:
            break
    
    # Calculate scale factors based on original dimensions
    # Original dimensions were likely 13.33" x 7.50" (16:9) = 12192000 x 6858000 EMU
    # If shapes are out of bounds, estimate original dimensions from current shape positions
    if shapes_need_scaling:
        print("\n⚠ Detected shapes extending beyond slide boundaries.")
        print("  Estimating original dimensions from shape positions...")
        
        # Find maximum extent of shapes across first few slides
        max_right = 0
        max_bottom = 0
        for i, slide in enumerate(slides_list[:min(5, len(slides_list))]):
            for shape in slide.shapes:
                max_right = max(max_right, shape.left + shape.width)
                max_bottom = max(max_bottom, shape.top + shape.height)
        
        # Estimate original dimensions (add some padding)
        estimated_original_width = max(max_right * 1.05, 12192000)  # At least 13.33"
        estimated_original_height = max(max_bottom * 1.05, 6858000)  # At least 7.50"
        
        print(f"  Estimated original dimensions: {estimated_original_width/914400:.2f}\" x {estimated_original_height/914400:.2f}\"")
        
        # Calculate scale factors from estimated original to target
        scale_x = SLIDE_WIDTH_EMU / estimated_original_width
        scale_y = SLIDE_HEIGHT_EMU / estimated_original_height
    else:
        # Normal case: calculate scale factors from current slide dimensions
        scale_x = SLIDE_WIDTH_EMU / current_width_emu
        scale_y = SLIDE_HEIGHT_EMU / current_height_emu
    
    print(f"\nScale factors: X={scale_x:.4f}, Y={scale_y:.4f}")
    
    # Set new slide dimensions FIRST (before scaling shapes)
    pres.slide_width = SLIDE_WIDTH_EMU
    pres.slide_height = SLIDE_HEIGHT_EMU
    
    print(f"\nResizing to: {SLIDE_WIDTH_IN:.2f}\" x {SLIDE_HEIGHT_IN:.2f}\"")
    print(f"Resizing to: {SLIDE_WIDTH_EMU} EMU x {SLIDE_HEIGHT_EMU} EMU")
    print(f"Resizing to: {SLIDE_WIDTH_PT}pt x {SLIDE_HEIGHT_PT}pt (16:9)")
    
    # Scale all shapes in all slides
    print(f"\nScaling {len(pres.slides)} slide(s)...")
    for slide_idx, slide in enumerate(pres.slides, 1):
        shape_count = 0
        for shape in slide.shapes:
            try:
                scale_shape(shape, scale_x, scale_y)
                shape_count += 1
            except Exception as e:
                print(f"  Warning: Could not scale shape on slide {slide_idx}: {e}")
        
        if slide_idx % 5 == 0 or slide_idx == len(pres.slides):
            print(f"  Processed slide {slide_idx}/{len(pres.slides)} ({shape_count} shapes)")
    
    # Save presentation
    print(f"\nSaving to: {output_path}")
    pres.save(str(output_path))
    
    print(f"✓ Successfully resized {len(pres.slides)} slide(s) with all elements scaled proportionally")
    
    return True


if __name__ == "__main__":
    if len(sys.argv) < 2:
        print("Usage: python resize_slide.py <input_pptx> [output_pptx]")
        print("\nExample:")
        print("  python resize_slide.py ref/Available _Slide.pptx")
        print("  python resize_slide.py ref/Available _Slide.pptx ref/Available _Slide_resized.pptx")
        sys.exit(1)
    
    input_pptx = sys.argv[1]
    output_pptx = sys.argv[2] if len(sys.argv) > 2 else None
    
    success = resize_presentation(input_pptx, output_pptx)
    sys.exit(0 if success else 1)

